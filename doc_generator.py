import os
import csv
import json
import mimetypes
import inspect
from datetime import datetime
from typing import Any, Dict, List, Optional
from urllib.parse import quote

from pydantic import Field

try:
    # pydantic v2
    from pydantic.fields import FieldInfo
except Exception:
    # pydantic v1 fallback
    FieldInfo = object  # type: ignore


class Tools:
    """
    Open WebUI Tools: gerar e analisar arquivos CSV, PDF, DOCX e PPTX.

    ✅ Atualizado para anexar arquivo via evento do Open WebUI:
      - Emite "files" e "chat:message:files" (compatibilidade)
      - Suporta __event_emitter__ síncrono OU assíncrono
      - Mantém fallback via download_url

    Env:
    - TOOLS_OUTPUT_DIR (default: /app/backend/tool_outputs)
    - PUBLIC_BASE_URL (default: "")           # opcional (para URL absoluta, se quiser)
    - DOWNLOADS_ROUTE_PATH (default: /downloads)
    """

    def __init__(self):
        self.output_dir = os.getenv("TOOLS_OUTPUT_DIR", "/app/backend/tool_outputs")
        self.public_base_url = (os.getenv("PUBLIC_BASE_URL", "") or "").strip().rstrip("/")

        drp = (os.getenv("DOWNLOADS_ROUTE_PATH", "/downloads") or "").strip()
        if not drp.startswith("/"):
            drp = "/" + drp
        drp = drp.rstrip("/") or "/downloads"
        self.downloads_route_path = drp

        os.makedirs(self.output_dir, exist_ok=True)

    # -----------------------------
    # Normalização (corrige FieldInfo)
    # -----------------------------
    def _is_fieldinfo(self, v: Any) -> bool:
        try:
            return isinstance(v, FieldInfo)
        except Exception:
            return False

    def _norm_str(self, v: Any, default: str = "") -> str:
        if self._is_fieldinfo(v) or v is None:
            return default
        return str(v)

    def _norm_list(self, v: Any, default: Optional[List[Any]] = None) -> List[Any]:
        if default is None:
            default = []
        if self._is_fieldinfo(v) or v is None:
            return default
        if isinstance(v, list):
            return v
        if isinstance(v, str):
            s = v.strip()
            if (s.startswith("[") and s.endswith("]")) or (s.startswith("{") and s.endswith("}")):
                try:
                    parsed = json.loads(s)
                    return parsed if isinstance(parsed, list) else default
                except Exception:
                    return default
        return default

    def _norm_int(self, v: Any, default: int) -> int:
        if self._is_fieldinfo(v) or v is None:
            return default
        try:
            return int(v)
        except Exception:
            return default

    # -----------------------------
    # Helpers
    # -----------------------------
    def _json(self, payload: Dict[str, Any]) -> str:
        return json.dumps(payload, ensure_ascii=False, indent=2)

    def _now_stamp(self) -> str:
        return datetime.now().strftime("%Y%m%d_%H%M%S")

    def _safe_basename(self, name: str, default: str = "arquivo") -> str:
        name = (name or "").strip()
        if not name:
            name = default

        name = name.replace("\\", "/")
        name = name.split("/")[-1]
        name = name.replace("..", "_").replace("\x00", "_")

        allowed = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789-_."
        cleaned = "".join(c if c in allowed else "_" for c in name).strip("._")
        return cleaned or default

    def _join_out(self, filename: str) -> str:
        path = os.path.abspath(os.path.join(self.output_dir, filename))
        base = os.path.abspath(self.output_dir)
        if not (path == base or path.startswith(base + os.sep)):
            raise ValueError("Caminho inválido (path traversal detectado).")
        return path

    def _file_info(self, path: str) -> Dict[str, Any]:
        st = os.stat(path)
        return {
            "path": os.path.abspath(path),
            "filename": os.path.basename(path),
            "size_bytes": st.st_size,
            "modified_at": datetime.fromtimestamp(st.st_mtime).isoformat(),
            "mime_type": mimetypes.guess_type(path)[0] or "application/octet-stream",
        }

    def _download_url_for(self, path: str) -> str:
        """
        Retorna sempre um caminho relativo /downloads/<arquivo> (com URL-encoding),
        e se PUBLIC_BASE_URL existir, retorna absoluto.
        """
        fname = os.path.basename(path)
        safe = quote(fname)
        rel = f"{self.downloads_route_path}/{safe}"
        return f"{self.public_base_url}{rel}" if self.public_base_url else rel

    async def _emit_any(self, emitter, payload: dict) -> None:
        """
        Emite evento suportando emitter async OU sync.
        """
        if not emitter:
            return

        # Alguns emitters são async (awaitable), outros sync.
        try:
            result = emitter(payload)
            if inspect.isawaitable(result):
                await result
        except Exception:
            # não quebra a tool
            return

    async def _emit_notification(self, __event_emitter__, content: str, level: str = "warning") -> None:
        """
        Notificação (toast) para debug opcional.
        """
        if not __event_emitter__ or not content:
            return
        await self._emit_any(
            __event_emitter__,
            {
                "type": "notification",
                "data": {"type": level, "content": content},
            },
        )

    async def _emit_file_attachment(
        self,
        __event_emitter__,
        *,
        name: str,
        url: str,
        debug: bool = False,
    ) -> Optional[str]:
        """
        Emite evento de files para o Open WebUI renderizar como anexo.

        Retorna string de erro (se debug=True), caso falhe.
        """
        if not __event_emitter__:
            return "no_event_emitter"
        if not name or not url:
            return "missing_name_or_url"

        file_obj = {"name": name, "url": url}

        # A doc lista 'files' e 'chat:message:files' como tipos válidos.
        # Emitimos os dois para compatibilidade. :contentReference[oaicite:3]{index=3}
        try:
            await self._emit_any(
                __event_emitter__,
                {"type": "files", "data": {"files": [file_obj]}},
            )
            await self._emit_any(
                __event_emitter__,
                {"type": "chat:message:files", "data": {"files": [file_obj]}},
            )
            return None
        except Exception as e:
            err = f"emit_failed: {e}"
            if debug:
                await self._emit_notification(__event_emitter__, f"Falha ao anexar arquivo: {err}", "warning")
            return err

    def _missing_dep(self, dep_name: str, pip_name: Optional[str] = None) -> str:
        pip_name = pip_name or dep_name
        return self._json(
            {
                "ok": False,
                "error": f"Dependência ausente: {dep_name}. Instale com: pip install {pip_name}",
            }
        )

    # -----------------------------
    # CSV
    # -----------------------------
    async def generate_csv(
        self,
        filename: str = Field(..., description="Nome do arquivo CSV (ex: dados.csv)"),
        headers: List[str] = Field(..., description="Lista de colunas"),
        rows: List[List[Any]] = Field(default_factory=list, description="Linhas"),
        delimiter: str = Field(";", description="Delimitador do CSV"),
        encoding: str = Field("utf-8", description="Encoding do arquivo"),
        __event_emitter__=None,
        __metadata__=None,
        debug_events: bool = Field(False, description="Se True, inclui info de debug de eventos no retorno"),
    ) -> str:
        try:
            filename = self._norm_str(filename, f"dados_{self._now_stamp()}.csv")
            headers = [self._norm_str(h) for h in self._norm_list(headers)]
            rows = self._norm_list(rows)

            delimiter = self._norm_str(delimiter, ";")
            encoding = self._norm_str(encoding, "utf-8")

            safe = self._safe_basename(filename, default=f"dados_{self._now_stamp()}.csv")
            if not safe.lower().endswith(".csv"):
                safe += ".csv"
            path = self._join_out(safe)

            for i, r in enumerate(rows):
                rr = r if isinstance(r, list) else []
                if len(rr) != len(headers):
                    return self._json(
                        {
                            "ok": False,
                            "error": f"Linha {i} tem {len(rr)} colunas, mas headers tem {len(headers)}.",
                        }
                    )

            with open(path, "w", newline="", encoding=encoding) as f:
                w = csv.writer(f, delimiter=delimiter)
                w.writerow(headers)
                w.writerows(rows)

            info = self._file_info(path)
            download_url = self._download_url_for(path)
            info["download_url"] = download_url
            info["action"] = "generate_csv"

            err = await self._emit_file_attachment(
                __event_emitter__, name=info["filename"], url=download_url, debug=debug_events
            )
            if debug_events and err:
                info["event_emit_error"] = err
                info["function_calling_mode_hint"] = (__metadata__ or {}).get("function_calling", "")

            return self._json({"ok": True, **info})
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})

    def analyze_csv(
        self,
        path: str = Field(..., description="Caminho do CSV para análise"),
        delimiter: Optional[str] = Field(None, description="Delimitador (ou None)"),
        encoding: str = Field("utf-8", description="Encoding do arquivo"),
        sample_rows: int = Field(10, description="Quantidade de linhas de amostra"),
    ) -> str:
        try:
            path = self._norm_str(path)
            encoding = self._norm_str(encoding, "utf-8")
            sample_rows = self._norm_int(sample_rows, 10)

            path = os.path.abspath(path)
            if not os.path.exists(path):
                return self._json({"ok": False, "error": "Arquivo não encontrado", "path": path})

            delim = delimiter
            if self._is_fieldinfo(delim):
                delim = None

            if delim is None:
                with open(path, "r", encoding=encoding, errors="ignore") as f:
                    head = f.read(4096)
                delim = ";" if head.count(";") >= head.count(",") else ","

            with open(path, "r", encoding=encoding, errors="ignore", newline="") as f:
                reader = csv.reader(f, delimiter=delim)
                all_rows = list(reader)

            if not all_rows:
                return self._json({"ok": True, "action": "analyze_csv", "empty": True, **self._file_info(path)})

            headers = all_rows[0]
            data = all_rows[1:]
            n_rows = len(data)
            n_cols = len(headers)

            sample = data[: max(0, sample_rows)]

            empty_counts = [0] * n_cols
            for r in data:
                for j in range(min(n_cols, len(r))):
                    if (r[j] is None) or (str(r[j]).strip() == ""):
                        empty_counts[j] += 1

            return self._json(
                {
                    "ok": True,
                    "action": "analyze_csv",
                    **self._file_info(path),
                    "delimiter": delim,
                    "columns": headers,
                    "num_columns": n_cols,
                    "num_data_rows": n_rows,
                    "empty_cells_per_column": dict(zip(headers, empty_counts)),
                    "sample_rows": sample,
                }
            )
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})

    # -----------------------------
    # PDF
    # -----------------------------
    async def generate_pdf(
        self,
        filename: str = Field(..., description="Nome do PDF (ex: relatorio.pdf)"),
        title: str = Field("Relatório", description="Título do documento"),
        paragraphs: List[str] = Field(default_factory=list, description="Parágrafos"),
        __event_emitter__=None,
        __metadata__=None,
        debug_events: bool = Field(False, description="Se True, inclui info de debug de eventos no retorno"),
    ) -> str:
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas
        except Exception:
            return self._missing_dep("reportlab")

        try:
            filename = self._norm_str(filename, f"relatorio_{self._now_stamp()}.pdf")
            title = self._norm_str(title, "Relatório")
            paragraphs = [self._norm_str(p) for p in self._norm_list(paragraphs)]

            safe = self._safe_basename(filename, default=f"relatorio_{self._now_stamp()}.pdf")
            if not safe.lower().endswith(".pdf"):
                safe += ".pdf"
            path = self._join_out(safe)

            c = canvas.Canvas(path, pagesize=A4)
            _, height = A4

            y = height - 72
            c.setFont("Helvetica-Bold", 16)
            c.drawString(72, y, title)
            y -= 36

            c.setFont("Helvetica", 11)
            for p in paragraphs:
                lines = str(p).splitlines() or [""]
                for line in lines:
                    if y < 72:
                        c.showPage()
                        y = height - 72
                        c.setFont("Helvetica", 11)
                    c.drawString(72, y, line[:1200])
                    y -= 14
                y -= 8

            c.save()

            info = self._file_info(path)
            download_url = self._download_url_for(path)
            info["download_url"] = download_url
            info["action"] = "generate_pdf"

            err = await self._emit_file_attachment(__event_emitter__, name=info["filename"], url=download_url, debug=debug_events)
            if debug_events and err:
                info["event_emit_error"] = err
                info["function_calling_mode_hint"] = (__metadata__ or {}).get("function_calling", "")

            return self._json({"ok": True, **info})
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})

    def analyze_pdf(
        self,
        path: str = Field(..., description="Caminho do PDF para análise"),
        max_pages: int = Field(10, description="Máximo de páginas"),
        max_chars: int = Field(4000, description="Máximo de caracteres"),
    ) -> str:
        path = os.path.abspath(self._norm_str(path))
        max_pages = self._norm_int(max_pages, 10)
        max_chars = self._norm_int(max_chars, 4000)

        if not os.path.exists(path):
            return self._json({"ok": False, "error": "Arquivo não encontrado", "path": path})

        try:
            from pypdf import PdfReader
        except Exception:
            try:
                from PyPDF2 import PdfReader  # type: ignore
            except Exception:
                return self._missing_dep("pypdf", "pypdf (ou PyPDF2)")

        try:
            reader = PdfReader(path)
            num_pages = len(reader.pages)
            pages_to_read = min(num_pages, max_pages)

            extracted = []
            for i in range(pages_to_read):
                try:
                    txt = reader.pages[i].extract_text() or ""
                except Exception:
                    txt = ""
                if txt.strip():
                    extracted.append(txt)

            full_text = "\n\n".join(extracted).strip()
            if len(full_text) > max_chars:
                full_text = full_text[:max_chars] + "\n...[truncado]"

            return self._json(
                {
                    "ok": True,
                    "action": "analyze_pdf",
                    **self._file_info(path),
                    "num_pages": num_pages,
                    "pages_read": pages_to_read,
                    "text_sample": full_text,
                }
            )
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})

    # -----------------------------
    # Word (DOCX)
    # -----------------------------
    async def generate_word(
        self,
        filename: str = Field(..., description="Nome do DOCX (ex: doc.docx)"),
        title: str = Field("Documento", description="Título do documento"),
        paragraphs: List[str] = Field(default_factory=list, description="Parágrafos"),
        __event_emitter__=None,
        __metadata__=None,
        debug_events: bool = Field(False, description="Se True, inclui info de debug de eventos no retorno"),
    ) -> str:
        try:
            from docx import Document
        except Exception:
            return self._missing_dep("python-docx")

        try:
            filename = self._norm_str(filename, f"documento_{self._now_stamp()}.docx")
            title = self._norm_str(title, "Documento")
            paragraphs = [self._norm_str(p) for p in self._norm_list(paragraphs)]

            safe = self._safe_basename(filename, default=f"documento_{self._now_stamp()}.docx")
            if not safe.lower().endswith(".docx"):
                safe += ".docx"
            path = self._join_out(safe)

            doc = Document()
            doc.add_heading(title, level=1)
            for p in paragraphs:
                doc.add_paragraph(p)
            doc.save(path)

            info = self._file_info(path)
            download_url = self._download_url_for(path)
            info["download_url"] = download_url
            info["action"] = "generate_word"

            err = await self._emit_file_attachment(__event_emitter__, name=info["filename"], url=download_url, debug=debug_events)
            if debug_events and err:
                info["event_emit_error"] = err
                info["function_calling_mode_hint"] = (__metadata__ or {}).get("function_calling", "")

            return self._json({"ok": True, **info})
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})

    def analyze_word(
        self,
        path: str = Field(..., description="Caminho do DOCX para análise"),
        max_paragraphs: int = Field(30, description="Máximo de parágrafos"),
        max_chars: int = Field(4000, description="Máximo de caracteres"),
    ) -> str:
        path = os.path.abspath(self._norm_str(path))
        max_paragraphs = self._norm_int(max_paragraphs, 30)
        max_chars = self._norm_int(max_chars, 4000)

        if not os.path.exists(path):
            return self._json({"ok": False, "error": "Arquivo não encontrado", "path": path})

        try:
            from docx import Document
        except Exception:
            return self._missing_dep("python-docx")

        try:
            doc = Document(path)
            paras = [p.text for p in doc.paragraphs if (p.text or "").strip()]
            sample = paras[:max_paragraphs]
            text_sample = "\n".join(sample).strip()
            if len(text_sample) > max_chars:
                text_sample = text_sample[:max_chars] + "\n...[truncado]"

            return self._json(
                {
                    "ok": True,
                    "action": "analyze_word",
                    **self._file_info(path),
                    "num_paragraphs": len(paras),
                    "paragraphs_read": len(sample),
                    "text_sample": text_sample,
                }
            )
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})

    # -----------------------------
    # PowerPoint (PPTX)
    # -----------------------------
    async def generate_ppt(
        self,
        filename: str = Field(..., description="Nome do PPTX (ex: deck.pptx)"),
        title: str = Field("Apresentação", description="Título da capa"),
        slides: List[Dict[str, Any]] = Field(
            default_factory=list,
            description="Slides: [{'title': str, 'bullets': [str], 'notes': str(opcional)}]",
        ),
        subtitle: str = Field("", description="Subtítulo opcional na capa"),
        author: str = Field("", description="Autor opcional na capa"),
        __event_emitter__=None,
        __metadata__=None,
        debug_events: bool = Field(False, description="Se True, inclui info de debug de eventos no retorno"),
    ) -> str:
        try:
            from pptx import Presentation
        except Exception:
            return self._missing_dep("python-pptx")

        try:
            filename = self._norm_str(filename, f"apresentacao_{self._now_stamp()}.pptx")
            title = self._norm_str(title, "Apresentação")
            subtitle = self._norm_str(subtitle, "")
            author = self._norm_str(author, "")
            slides = self._norm_list(slides, default=[])

            safe = self._safe_basename(filename, default=f"apresentacao_{self._now_stamp()}.pptx")
            if not safe.lower().endswith(".pptx"):
                safe += ".pptx"
            path = self._join_out(safe)

            prs = Presentation()

            # Capa
            cover_layout = prs.slide_layouts[0]
            cover = prs.slides.add_slide(cover_layout)
            cover.shapes.title.text = title

            if len(cover.placeholders) > 1:
                meta: List[str] = []
                if subtitle:
                    meta.append(subtitle)
                meta.append(datetime.now().strftime("%Y-%m-%d"))
                if author:
                    meta.append(author)
                cover.placeholders[1].text = " • ".join([m for m in meta if m])

            # Conteúdo
            content_layout = prs.slide_layouts[1]
            for item in slides:
                if not isinstance(item, dict):
                    continue
                st = self._norm_str(item.get("title", ""), "")
                bullets = item.get("bullets", [])
                notes = item.get("notes", None)

                bullets_list = [self._norm_str(b) for b in self._norm_list(bullets)]
                notes_str = self._norm_str(notes, "") if notes is not None else ""

                s = prs.slides.add_slide(content_layout)
                s.shapes.title.text = st

                tf = s.shapes.placeholders[1].text_frame
                tf.clear()

                for i, b in enumerate(bullets_list):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = b

                if notes_str:
                    try:
                        s.notes_slide.notes_text_frame.text = notes_str
                    except Exception:
                        pass

            prs.save(path)

            info = self._file_info(path)
            download_url = self._download_url_for(path)
            info["download_url"] = download_url
            info["action"] = "generate_ppt"

            err = await self._emit_file_attachment(__event_emitter__, name=info["filename"], url=download_url, debug=debug_events)
            if debug_events and err:
                info["event_emit_error"] = err
                info["function_calling_mode_hint"] = (__metadata__ or {}).get("function_calling", "")

            return self._json({"ok": True, **info})
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})

    def analyze_ppt(
        self,
        path: str = Field(..., description="Caminho do PPTX para análise"),
        max_slides: int = Field(20, description="Máximo de slides"),
        max_chars: int = Field(4000, description="Máximo de caracteres"),
    ) -> str:
        path = os.path.abspath(self._norm_str(path))
        max_slides = self._norm_int(max_slides, 20)
        max_chars = self._norm_int(max_chars, 4000)

        if not os.path.exists(path):
            return self._json({"ok": False, "error": "Arquivo não encontrado", "path": path})

        try:
            from pptx import Presentation
        except Exception:
            return self._missing_dep("python-pptx")

        try:
            prs = Presentation(path)
            num_slides = len(prs.slides)
            slides_to_read = min(num_slides, max_slides)

            extracted: List[Dict[str, Any]] = []
            for i in range(slides_to_read):
                slide = prs.slides[i]
                title = ""
                if slide.shapes.title is not None:
                    try:
                        title = slide.shapes.title.text or ""
                    except Exception:
                        title = ""

                texts: List[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                        t = (shape.text_frame.text or "").strip()
                        if t:
                            texts.append(t)

                extracted.append({"slide_index": i, "title": title, "text_blocks": texts[:20]})

            joined: List[str] = []
            for s in extracted:
                if s.get("title"):
                    joined.append(s["title"])
                for tb in s.get("text_blocks", []):
                    joined.append(tb)

            text_sample = "\n\n".join(joined).strip()
            if len(text_sample) > max_chars:
                text_sample = text_sample[:max_chars] + "\n...[truncado]"

            return self._json(
                {
                    "ok": True,
                    "action": "analyze_ppt",
                    **self._file_info(path),
                    "num_slides": num_slides,
                    "slides_read": slides_to_read,
                    "slides_summary": extracted,
                    "text_sample": text_sample,
                }
            )
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})

    # -----------------------------
    # Analyzer universal
    # -----------------------------
    def analyze_file(
        self,
        path: str = Field(..., description="Caminho do arquivo para análise (csv/pdf/docx/pptx)"),
    ) -> str:
        try:
            path_abs = os.path.abspath(self._norm_str(path))
            ext = os.path.splitext(path_abs)[1].lower()

            if ext == ".csv":
                # chamada sync de analyze_csv
                return self.analyze_csv(path_abs)
            if ext == ".pdf":
                return self.analyze_pdf(path_abs)
            if ext == ".docx":
                return self.analyze_word(path_abs)
            if ext == ".pptx":
                return self.analyze_ppt(path_abs)

            return self._json(
                {
                    "ok": False,
                    "error": f"Extensão não suportada: {ext}. Suportadas: .csv .pdf .docx .pptx",
                    "path": path_abs,
                }
            )
        except Exception as e:
            return self._json({"ok": False, "error": str(e)})
